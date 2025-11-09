import os
import uuid
import logging
import threading
from contextlib import asynccontextmanager
from datetime import date
from http.client import HTTPException

from io import BytesIO
from pypdf import PdfReader
from pptx import Presentation

from dotenv import load_dotenv
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough, RunnableLambda

from pydantic import BaseModel
from fastapi.responses import JSONResponse
from fastapi import FastAPI, UploadFile, File
from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

load_dotenv()
os.makedirs("logs", exist_ok=True)

# Setup logging
log_file_name = f"logs/{date.today()}_backend.log"
logging.basicConfig(
    level=logging.INFO,
    format="[%(asctime)s] [%(levelname)s] [%(name)s] %(message)s",
    handlers=[
        logging.FileHandler(log_file_name, encoding="utf-8")
    ]
)

logger = logging.getLogger("API Logger")

# Constants
TASKS = {}
PERSIST_DIR = "./Chroma"
COLLECTION_NAME = "vector_db"

class ChatRequest(BaseModel):
    message: str


def log_prompt(prompt):
    """Utility function to log prompts"""
    logger.info(f"Generated prompt for the question:\n{prompt}")
    return prompt

def extract_text_from_pdf(file_bytes: bytes) -> str:
    parts = []
    reader = PdfReader(BytesIO(file_bytes))

    for page in reader.pages:
        text = page.extract_text() or ""
        parts.append(text)

    return "\n".join(parts)

def extract_text_from_ppt(file_bytes: bytes) -> str:
    parts = []
    reader = Presentation(BytesIO(file_bytes))

    for slide in reader.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)

    return "\n".join(parts)

def process_file(task_id: str, data: str, file_name: str = "upload"):
    """
    Background task to chunk uploaded file and save it in vector store

    Args:
        task_id: internal unique task identifier
        file_bytes: bytes of the uploaded file
        file_name: name of the uploaded file

    Returns:
        Nothing, processes the file and updates task_id status internally
    """
    try:
        TASKS[task_id] = {"status": "processing"}
        logger.info(f"Started processing task id: {task_id}, status: processing")
        file = Document(page_content=data, metadata={"source": file_name})

        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        documents = splitter.split_documents([file])

        app.state.vector_store.add_documents(documents)
        app.state.vector_store.persist()
        TASKS[task_id] = {"status": "success", "chunks": len(documents)}
        logger.info(f"Task id {task_id} completed, status: success")
        return True
    except Exception as e:
        TASKS[task_id] = {"status": "failed", "error": str(e)}
        logger.error(f"Task id {task_id} failed, error\n{e}, status: fail")


@asynccontextmanager
async def lifespan(app: FastAPI):
    llm = ChatOpenAI(model="gpt-4o-mini",
                     api_key=os.environ["OPENAI_API_KEY"])
    embeddings = OpenAIEmbeddings(model="text-embedding-3-small",
                                  api_key=os.environ["OPENAI_API_KEY"])
    vector_store = Chroma(collection_name=COLLECTION_NAME,
                          embedding_function=embeddings,
                          persist_directory=PERSIST_DIR)

    app.state.llm = llm
    app.state.embeddings = embeddings
    app.state.vector_store = vector_store
    logger.info("Loaded embeddings model and vector store")
    yield

app = FastAPI(title="Docs Agent", version="0.0.0.0", lifespan=lifespan)

@app.post("/upload")
async def file_upload(file: UploadFile = File(...)):
    file_bytes = await file.read()
    file_name = file.filename
    task_id = str(uuid.uuid4())
    TASKS[task_id] = {"status": "accepted"}
    logger.info(f"Accepted task id: {task_id}, status: accepted")

    ext = os.path.splitext(file_name)[1].lower()

    if ext == ".pdf":
        data = extract_text_from_pdf(file_bytes)
    elif ext in (".ppt", ".pptx"):
        data = extract_text_from_ppt(file_bytes)
    else:
        data = file_bytes.decode("utf-8", errors="ignore")

    thread = threading.Thread(target=process_file, args=(task_id, data, file_name))
    thread.start()

    return JSONResponse(
        status_code=202,
        content={"task_id": task_id, "status": "accepted"}
    )

@app.get("/status/{task_id}")
def get_status(task_id: str):
    task = TASKS.get(task_id)
    logger.info(f"Looking for task id {task_id}")

    if not task:
        return {"status": "Not Found"}

    logger.info(f"Returned task {task}")
    return task

@app.post("/chat")
def chat(req: ChatRequest):
    question = req.message.strip()

    if not question:
        raise HTTPException(status_code=400, detail="Empty message")

    logger.info(f"Answering question {question}")

    # Get chain components
    llm = app.state.llm
    vector_store = app.state.vector_store

    # Create a retriever
    retriever = vector_store.as_retriever(search_kwargs={"k": 3})

    # Create Prompt
    template = ChatPromptTemplate.from_template("""
    You are a meeting assistant agent. Answer the user's questions with the help of given meeting files.
    
    Context:
    {context}
    
    Question:
    {question}
    
    Answer:
    """)

    # Create chain
    chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | template
        | RunnableLambda(log_prompt)
        | llm
        | StrOutputParser()
    )

    try:
        answer = chain.invoke(question)
        logger.info(f"Got answer {answer}")
        return {"answer": answer}
    except Exception as e:
        logger.error(f"Error getting answer, {e}")

    return {"answer": None}
